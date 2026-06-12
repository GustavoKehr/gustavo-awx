"""
Gera apresentacao PowerPoint da Stack de Observabilidade.
Execute: py -3 gerar_apresentacao.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy

# Paleta de cores
BG_DARK    = RGBColor(0x0F, 0x11, 0x17)
BG_CARD    = RGBColor(0x1E, 0x23, 0x30)
ACCENT     = RGBColor(0x3B, 0x82, 0xF6)   # azul
ACCENT2    = RGBColor(0x22, 0xD3, 0xEE)   # ciano  (Loki)
ACCENT3    = RGBColor(0xF5, 0x9E, 0x0B)   # laranja (Prometheus)
ACCENT4    = RGBColor(0x10, 0xB9, 0x81)   # verde   (Alloy)
ACCENT5    = RGBColor(0xA7, 0x8B, 0xFA)   # roxo    (Grafana)
RED        = RGBColor(0xEF, 0x44, 0x44)
YELLOW     = RGBColor(0xF5, 0x9E, 0x0B)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0x94, 0xA3, 0xB8)
GRAY_LIGHT = RGBColor(0xE2, 0xE8, 0xF0)
GREEN      = RGBColor(0x10, 0xB9, 0x81)
BG_CODE    = RGBColor(0x0D, 0x11, 0x1A)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide_add(title_txt=""):
    slide = prs.slides.add_slide(BLANK)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_DARK
    return slide


def add_rect(slide, l, t, w, h, fill_color):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h, size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox


def accent_bar(slide, color=ACCENT):
    add_rect(slide, 0, 0, 13.33, 0.08, color)


def slide_title(title, subtitle="", accent_color=ACCENT):
    slide = slide_add(title)
    accent_bar(slide, accent_color)
    add_text(slide, title, 0.8, 0.25, 11.5, 1.0, size=36, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.8, 1.2, 11.5, 0.6, size=16, color=GRAY)
    return slide


def badge(slide, text, l, t, w=1.8, h=0.35, color=ACCENT):
    add_rect(slide, l, t, w, h, color)
    add_text(slide, text, l+0.05, t+0.02, w-0.1, h-0.05,
             size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 1 - Capa
# =============================================================================
s = slide_add("Capa")
add_rect(s, 0, 0, 13.33, 7.5, BG_DARK)
add_rect(s, 0, 0, 0.5, 7.5, ACCENT)
add_rect(s, 0.5, 3.2, 12.83, 0.06, ACCENT)

add_text(s, "Stack de Observabilidade", 1.0, 0.7, 11.3, 1.5, size=44, bold=True, color=WHITE)
add_text(s, "Loki  |  Grafana  |  Prometheus  |  Alloy", 1.0, 2.1, 11.3, 0.8, size=22, color=ACCENT2)
add_text(s, "Monitoramento unificado de logs e metricas\npara ambientes corporativos",
         1.0, 3.4, 10.0, 1.0, size=16, color=GRAY)

for i, (label, color) in enumerate([("Loki 3.6", ACCENT2), ("Grafana 13.0", ACCENT5),
                                      ("Prometheus 3.11", ACCENT3), ("Alloy 1.9", ACCENT4)]):
    badge(s, label, 1.0 + i * 2.5, 4.6, 2.2, 0.38, color)

add_text(s, "Infraestrutura  |  2026", 1.0, 6.7, 11.3, 0.5, size=12, color=GRAY, italic=True)


# =============================================================================
# SLIDE 2 - O Problema
# =============================================================================
s = slide_title("O Problema", "Como esta hoje - sem observabilidade", ACCENT3)

boxes = [
    ("Logs espalhados",        "Cada servidor tem seus logs em /var/log.\nNao ha busca centralizada.",          ACCENT3),
    ("Resposta lenta",         "Para investigar um erro, voce acessa\nservidor por servidor via SSH.",          ACCENT3),
    ("Sem visibilidade",       "CPU alta? Disco cheio? So descobre\nquando o usuario reclama.",                ACCENT3),
    ("Sem alertas proativos",  "Sem monitoramento, voce reage\naos problemas - nao os previne.",               ACCENT3),
]
for i, (title, body, color) in enumerate(boxes):
    x = 0.4 + i * 3.2
    add_rect(s, x, 1.9, 3.0, 4.2, BG_CARD)
    add_rect(s, x, 1.9, 3.0, 0.06, color)
    add_text(s, title, x+0.15, 2.05, 2.7, 0.6, size=13, bold=True, color=WHITE)
    add_text(s, body,  x+0.15, 2.75, 2.7, 3.0, size=11, color=GRAY)


# =============================================================================
# SLIDE 3 - A Solucao
# =============================================================================
s = slide_title("A Solucao", "Uma plataforma unica para logs, metricas e alertas", ACCENT)
add_text(s, "Stack open-source de mercado. Deploy automatizado via Ansible. Funciona sem internet nos servidores.",
         0.8, 1.6, 11.5, 0.6, size=13, color=GRAY_LIGHT)

comps = [
    ("Grafana",     "Interface web unificada\nDashboards + alertas\nVersao 13.0",        ACCENT5),
    ("Loki",        "Banco de dados de logs\nIndexacao por labels\nVersao 3.6",           ACCENT2),
    ("Prometheus",  "Banco de dados de metricas\nSeries temporais\nVersao 3.11",          ACCENT3),
    ("Alloy",       "Agente unificado\nColeta logs + metricas\nVersao 1.9",               ACCENT4),
]
for i, (name, desc, color) in enumerate(comps):
    x = 0.5 + i * 3.2
    add_rect(s, x, 2.4, 3.0, 3.9, BG_CARD)
    add_rect(s, x, 2.4, 3.0, 0.5, color)
    add_text(s, name, x+0.15, 2.45, 2.7, 0.45, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, desc, x+0.15, 3.05, 2.7, 2.8, size=11, color=GRAY_LIGHT)


# =============================================================================
# SLIDE 4 - Arquitetura
# =============================================================================
s = slide_title("Arquitetura", "Como os componentes se conectam", ACCENT4)

# Servidor central
add_rect(s, 5.2, 1.7, 4.2, 4.8, BG_CARD)
add_rect(s, 5.2, 1.7, 4.2, 0.06, ACCENT)
add_text(s, "obs-server :200", 5.35, 1.8, 3.9, 0.4, size=12, bold=True, color=ACCENT)
add_text(s, "Grafana :3000",     5.35, 2.35, 3.9, 0.38, size=11, bold=True, color=ACCENT5)
add_text(s, "Loki :3100",        5.35, 2.82, 3.9, 0.38, size=11, bold=True, color=ACCENT2)
add_text(s, "Prometheus :9090",  5.35, 3.29, 3.9, 0.38, size=11, bold=True, color=ACCENT3)
add_rect(s, 5.35, 3.75, 3.9, 0.03, RGBColor(0x2D, 0x37, 0x48))
add_text(s, "retencao 30 dias", 5.35, 3.85, 3.9, 0.35, size=9, color=GRAY, italic=True)
add_text(s, "prometheus remote-write", 5.35, 4.2, 3.9, 0.35, size=9, color=GRAY, italic=True)

# Agentes
for name, ip, y in [("obs-agent1", ":201", 1.9), ("obs-agent2", ":202", 3.9)]:
    add_rect(s, 0.8, y, 3.2, 1.5, BG_CARD)
    add_rect(s, 0.8, y, 3.2, 0.06, ACCENT4)
    add_text(s, f"{name} {ip}", 0.95, y+0.12, 2.9, 0.42, size=11, bold=True, color=ACCENT4)
    add_text(s, "Alloy :12345",   0.95, y+0.6,  2.9, 0.3,  size=10, color=GRAY)
    add_text(s, "logs + metricas", 0.95, y+0.95, 2.9, 0.3,  size=10, color=GRAY)

add_text(s, "push logs -> :3100\npush metrics -> :9090", 3.85, 2.45, 1.5, 0.7, size=9, color=GRAY, italic=True)
add_text(s, "push logs -> :3100\npush metrics -> :9090", 3.85, 4.25, 1.5, 0.7, size=9, color=GRAY, italic=True)

# Usuario
add_rect(s, 10.4, 2.8, 2.4, 1.3, BG_CARD)
add_text(s, "Time / NOC",    10.55, 2.9,  2.1, 0.42, size=11, bold=True, color=WHITE)
add_text(s, "browser :3000", 10.55, 3.38, 2.1, 0.4,  size=10, color=GRAY)
add_text(s, "query <- :3000", 9.65, 3.1, 0.9, 0.5, size=9, color=GRAY, italic=True)


# =============================================================================
# SLIDE 5 - Loki: Logs
# =============================================================================
s = slide_title("Loki - Centralizacao de Logs", "Versao 3.6 | Modo monolitico | TSDB schema v13", ACCENT2)

add_text(s, "O que e coletado automaticamente em cada servidor:", 0.8, 1.55, 11.5, 0.45, size=13, color=GRAY_LIGHT)

sources = [
    ("systemd journal",   "Todos os logs do SO:\nSSH, sudo, cron, servicos",   ACCENT2),
    ("/var/log/messages", "Mensagens gerais\ndo sistema",                       ACCENT2),
    ("/var/log/secure",   "Acessos, autenticacoes,\nlogins SSH",               ACCENT2),
    ("/var/log/cron",     "Execucoes agendadas\ncrontab",                      ACCENT2),
    ("Labels coletados",  "job, host, level\nFiltragem no Grafana",            ACCENT),
]
for i, (src, desc, color) in enumerate(sources):
    row = i // 3
    col = i % 3
    x = 0.8 + col * 4.1
    y = 2.1 + row * 1.9
    add_rect(s, x, y, 3.8, 1.7, BG_CARD)
    add_rect(s, x, y, 0.06, 1.7, color)
    add_text(s, src,  x+0.2, y+0.15, 3.4, 0.5, size=12, bold=True, color=WHITE)
    add_text(s, desc, x+0.2, y+0.7,  3.4, 0.8, size=10, color=GRAY)

add_text(s, "Exemplo LogQL:", 0.8, 5.9, 2.5, 0.38, size=11, bold=True, color=GRAY)
add_rect(s, 0.8, 6.3, 11.5, 0.85, BG_CODE)
add_text(s, '{job="systemd-journal"} |= "Failed password"  ->  todos os logins SSH recusados',
         1.0, 6.42, 11.0, 0.55, size=11, color=ACCENT2)


# =============================================================================
# SLIDE 6 - Prometheus: Metricas
# =============================================================================
s = slide_title("Prometheus - Metricas de Infraestrutura", "Versao 3.11 | Remote-write receiver ativado", ACCENT3)

add_text(s, "Metricas coletadas via Grafana Alloy (node_exporter integrado):",
         0.8, 1.55, 11.5, 0.45, size=13, color=GRAY_LIGHT)

metrics = [
    ("CPU",       "Uso por core e modo\n(user, system, idle, iowait)",  ACCENT3),
    ("Memoria",   "Total, disponivel,\nbuffers, cache, swap",            ACCENT3),
    ("Disco",     "Espaco usado/livre\npor particao + IOPS + latencia",  ACCENT3),
    ("Rede",      "Bytes TX/RX, pacotes,\nerros por interface",          ACCENT3),
    ("Sistema",   "Uptime, load average,\nprocessos, file descriptors",  ACCENT3),
    ("TCP/IP",    "Conexoes estabelecidas,\ntime_wait, sockets em uso",  ACCENT3),
]
for i, (name, desc, color) in enumerate(metrics):
    row = i // 3
    col = i % 3
    x = 0.8 + col * 4.1
    y = 2.1 + row * 1.9
    add_rect(s, x, y, 3.8, 1.7, BG_CARD)
    add_rect(s, x, y, 0.06, 1.7, color)
    add_text(s, name, x+0.2, y+0.15, 3.4, 0.5, size=13, bold=True, color=WHITE)
    add_text(s, desc, x+0.2, y+0.7,  3.4, 0.8, size=10, color=GRAY)

add_text(s, "Exemplo PromQL:", 0.8, 5.9, 2.5, 0.38, size=11, bold=True, color=GRAY)
add_rect(s, 0.8, 6.3, 11.5, 0.85, BG_CODE)
add_text(s, "100 - (avg by (host)(rate(node_cpu_seconds_total{mode='idle'}[5m]))*100)  ->  % CPU por host",
         1.0, 6.42, 11.0, 0.55, size=11, color=ACCENT3)


# =============================================================================
# SLIDE 7 - Grafana: Os 3 Dashboards
# =============================================================================
s = slide_title("Grafana - Dashboards Deployados", "3 dashboards prontos e em producao | Versao 13.0", ACCENT5)

add_text(s, "Dashboards criados programaticamente via Python e provisionados automaticamente pelo Ansible:",
         0.8, 1.55, 11.5, 0.45, size=13, color=GRAY_LIGHT)

dashboards = [
    (
        "Linux System Overview",
        "linux-overview-v1",
        "CPU | Memoria | Disco\nRede | Load | Uptime\nFile descriptors | TCP",
        "Prometheus",
        ACCENT3,
    ),
    (
        "Logs e Analise",
        "logs-analysis-v1",
        "Log rate por host\nErros + warnings (1h)\nSSH logins | Eventos sudo\nFalhas de servico systemd",
        "Loki",
        ACCENT2,
    ),
    (
        "Infrastructure Comparison",
        "infra-compare-v1",
        "CPU % todos os hosts\nRAM todos os hosts\nDisco | I/O | Rede\nComparacao lado a lado",
        "Prometheus + Loki",
        ACCENT,
    ),
]
for i, (title, uid, desc, source, color) in enumerate(dashboards):
    x = 0.5 + i * 4.25
    add_rect(s, x, 2.1, 3.9, 4.8, BG_CARD)
    add_rect(s, x, 2.1, 3.9, 0.5, color)
    add_text(s, title,  x+0.15, 2.13, 3.6, 0.45, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, desc,   x+0.15, 2.75, 3.6, 2.5,  size=10, color=GRAY_LIGHT)
    add_rect(s, x+0.15, 5.5, 3.6, 0.03, RGBColor(0x2D, 0x37, 0x48))
    add_text(s, f"uid: {uid}", x+0.15, 5.6, 3.6, 0.35, size=9, color=GRAY, italic=True)
    badge(s, source, x+0.15, 6.05, 3.6, 0.35, color)


# =============================================================================
# SLIDE 8 - Dashboards em Detalhe
# =============================================================================
s = slide_title("Dashboards - Paineis Principais", "O que voce ve em cada dashboard", ACCENT5)

rows = [
    ("Linux Overview",        "CPU breakdown por modo",      "Load average 1/5/15m",    "RAM com buffers/cache",    "Disco: espaco + IOPS + latencia",   "Rede: RX/TX + erros + TCP"),
    ("Logs e Analise",        "Total de linhas (1h)",        "Contagem de erros (1h)",  "Taxa de logs por host",    "SSH: logins ok vs falhas",          "Tail de erros em tempo real"),
    ("Infra Comparison",      "CPU % todos os hosts",        "RAM % todos os hosts",    "Disco / por host",        "I/O total por host",                "Taxa de erros (Loki)"),
]

y_start = 1.8
for ri, (dash, *paineis) in enumerate(rows):
    y = y_start + ri * 1.65
    add_rect(s, 0.4, y, 2.1, 1.4, BG_CARD)
    color = [ACCENT3, ACCENT2, ACCENT][ri]
    add_rect(s, 0.4, y, 2.1, 0.06, color)
    add_text(s, dash, 0.5, y+0.12, 1.9, 1.1, size=10, bold=True, color=color)

    for ci, painel in enumerate(paineis):
        px = 2.65 + ci * 1.78
        add_rect(s, px, y+0.1, 1.65, 1.2, RGBColor(0x16, 0x1B, 0x27))
        add_rect(s, px, y+0.1, 1.65, 0.04, color)
        add_text(s, painel, px+0.08, y+0.2, 1.5, 0.9, size=9, color=GRAY_LIGHT)

add_text(s, "Todos os dashboards provisionados em /var/lib/grafana/dashboards/ - carregados automaticamente pelo Grafana",
         0.8, 6.85, 11.5, 0.45, size=10, color=GRAY, italic=True)


# =============================================================================
# SLIDE 9 - Alertas
# =============================================================================
s = slide_title("Alertas Proativos", "Grafana Unified Alerting", ACCENT3)

add_text(s, "Regras de alerta que disparam automaticamente quando algo foge do normal:",
         0.8, 1.55, 11.5, 0.45, size=13, color=GRAY_LIGHT)

alerts = [
    ("CPU Alta",       "CPU > 80% por 5 minutos continuamente",            "critical"),
    ("Memoria",        "Memoria disponivel < 10% por 2 minutos",           "warning"),
    ("Disco Cheio",    "Disco / > 85% por 10 minutos",                     "critical"),
    ("Host Down",      "Agente parou de enviar dados por 5 minutos",       "warning"),
    ("Muitos Erros",   "Mais de 50 erros nos logs em 5 minutos",           "warning"),
]
for i, (name, cond, sev) in enumerate(alerts):
    y = 2.15 + i * 0.9
    color = RED if sev == "critical" else YELLOW
    add_rect(s, 0.8, y, 0.08, 0.72, color)
    add_text(s, name, 1.1, y+0.06, 2.8, 0.38, size=12, bold=True, color=WHITE)
    add_text(s, cond, 1.1, y+0.42, 4.0, 0.28, size=10, color=GRAY)
    add_rect(s, 5.2, y, 7.8, 0.72, BG_CARD)
    lbl = "Critical" if sev == "critical" else "Warning"
    add_text(s, lbl, 5.4, y+0.18, 1.4, 0.34, size=10, bold=True, color=color)
    add_text(s, "->  Email  |  Slack  |  Teams  |  Webhook",
             7.0, y+0.18, 5.8, 0.34, size=10, color=GRAY)

add_text(s, "Os alertas chegam no canal que voce definir: Slack, Teams, e-mail, ou qualquer webhook",
         0.8, 6.75, 11.5, 0.45, size=11, color=GRAY, italic=True)


# =============================================================================
# SLIDE 10 - Deploy / Implementacao
# =============================================================================
s = slide_title("Implementacao", "Ansible Air-Gapped | Sem internet nos servidores alvo", ACCENT4)

add_text(s, "Toda a instalacao e automatizada via Ansible. Os binarios ficam no controller - nenhum host precisa de acesso externo.",
         0.8, 1.55, 11.5, 0.5, size=13, color=GRAY_LIGHT)

steps = [
    ("1", "Baixar artefatos",        "Na maquina com internet: Loki, Grafana RPM, Prometheus, Alloy RPM (~520 MB total).",   ACCENT),
    ("2", "Transferir ao controller","Copiar pacote para o servidor Ansible via SCP. Sem internet necessaria apos isso.",     ACCENT4),
    ("3", "Configurar inventario",   "Editar hosts.yml com IPs + grupo observability_server (servidor) e linux_agents.",      ACCENT3),
    ("4", "Deploy do servidor",      "ansible-playbook site.yml --tags server  ->  instala Loki + Grafana + Prometheus.",     ACCENT5),
    ("5", "Deploy dos agentes",      "ansible-playbook 20_linux_agents.yml  ->  instala Alloy em todos os linux_agents.",    ACCENT2),
    ("6", "Novos servidores",        "Para adicionar agente futuro: so adicionar ao inventario e rodar -l novo-host.",        ACCENT4),
]
for i, (num, title, desc, color) in enumerate(steps):
    y = 2.1 + i * 0.86
    add_rect(s, 0.8, y, 0.55, 0.72, color)
    add_text(s, num, 0.8, y+0.04, 0.55, 0.64, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, title, 1.6, y+0.04, 3.2, 0.38, size=11, bold=True, color=WHITE)
    add_text(s, desc,  1.6, y+0.4,  11.2, 0.3,  size=9,  color=GRAY)


# =============================================================================
# SLIDE 11 - Licoes Aprendidas (instalacao manual)
# =============================================================================
s = slide_title("Licoes Aprendidas", "Problemas encontrados na instalacao manual air-gapped", ACCENT3)

add_text(s, "Validado com instalacao manual completa em AlmaLinux 9.8. Todos documentados no playbook.",
         0.8, 1.55, 11.5, 0.45, size=13, color=GRAY_LIGHT)

lessons = [
    (
        "SELinux: user_tmp_t",
        "Binarios copiados via SCP herdam contexto user_tmp_t.\nSystemd recusa executar -> status=203/EXEC.",
        "sudo restorecon -v /usr/local/bin/loki\nsudo restorecon -v /usr/local/bin/prometheus",
        RED,
    ),
    (
        "Loki 3.6: Consul KV",
        "Modo monolitico tenta usar Consul por padrao.\nSem Consul: failed to start store (localhost:8500).",
        "common:\n  ring:\n    kvstore:\n      store: inmemory",
        ACCENT2,
    ),
    (
        "Permissao logs: chown",
        "chmod 640 sozinho nao basta se grupo for root.\nGrupo precisa ser adm para alloy ter acesso.",
        "sudo chown root:adm /var/log/messages\nsudo chown root:adm /var/log/secure",
        ACCENT3,
    ),
    (
        "sudo PATH restrito",
        "/usr/local/bin fora do PATH seguro do sudo.\npromtool nao encontrado com sudo.",
        "sudo /usr/local/bin/promtool check config\n/etc/prometheus/prometheus.yml",
        ACCENT,
    ),
]
for i, (title, prob, fix, color) in enumerate(lessons):
    row = i // 2
    col = i % 2
    x = 0.5 + col * 6.4
    y = 2.1 + row * 2.3
    add_rect(s, x, y, 6.1, 2.1, BG_CARD)
    add_rect(s, x, y, 6.1, 0.06, color)
    add_text(s, title, x+0.15, y+0.12, 5.8, 0.42, size=12, bold=True, color=color)
    add_text(s, prob,  x+0.15, y+0.6,  5.8, 0.65, size=9,  color=GRAY)
    add_rect(s, x+0.15, y+1.28, 5.8, 0.65, BG_CODE)
    add_text(s, fix,   x+0.25, y+1.33, 5.6, 0.55, size=8,  color=ACCENT4)


# =============================================================================
# SLIDE 12 - Beneficios
# =============================================================================
s = slide_title("Beneficios", "", ACCENT)

benefits = [
    ("Resposta mais rapida",    "Identifique a causa raiz de incidentes\nem minutos, nao em horas.",              ACCENT),
    ("Alertas proativos",       "Seja notificado antes que\no usuario perceba o problema.",                       ACCENT4),
    ("Seguranca e auditoria",   "Todos os logins, comandos sudo\ne acessos SSH ficam registrados.",               ACCENT2),
    ("Visibilidade completa",   "CPU, memoria, disco, rede\nde todos os servidores em um painel.",               ACCENT3),
    ("Open-source",             "Grafana, Loki e Prometheus sao\nopen-source. Sem custo de licenca.",            ACCENT5),
    ("Deploy automatizado",     "Novo servidor monitorado\nem menos de 5 minutos via Ansible.",                  ACCENT4),
]
for i, (title, desc, color) in enumerate(benefits):
    row = i // 2
    col = i % 2
    x = 0.8 + col * 6.3
    y = 1.5 + row * 1.8
    add_rect(s, x, y, 6.0, 1.6, BG_CARD)
    add_rect(s, x, y, 6.0, 0.06, color)
    add_text(s, title, x+0.2, y+0.15, 5.6, 0.5, size=13, bold=True, color=WHITE)
    add_text(s, desc,  x+0.2, y+0.7,  5.6, 0.7, size=11, color=GRAY)


# =============================================================================
# SLIDE 13 - Proximos Passos
# =============================================================================
s = slide_title("Proximos Passos", "", ACCENT5)

steps2 = [
    ("Fase 1 - Infraestrutura base",     "Stack deployada: Loki + Grafana + Prometheus no servidor central.\nAgentes Alloy em obs-agent1 e obs-agent2. Dashboards provisionados.",          "CONCLUIDO",  GREEN),
    ("Fase 2 - Expandir agentes",        "Adicionar Alloy nos servidores de aplicacao, banco de dados e DMZ.\nCada novo host: 1 linha no inventario + 1 comando Ansible.",                 "Proximo",    ACCENT),
    ("Fase 3 - Alertas e notificacoes",  "Configurar regras de alerta para CPU, disco, host down e erros.\nIntegrar com Slack ou Teams para notificacao em tempo real.",                  "~2 dias",    ACCENT3),
    ("Fase 4 - Logs de aplicacoes",      "Adicionar paths de log das aplicacoes no config do Alloy.\nParsing de logs JSON para facilitar busca e correlacao.",                            "Por demanda", ACCENT5),
]
for i, (title, desc, effort, color) in enumerate(steps2):
    y = 1.7 + i * 1.35
    add_rect(s, 0.8, y, 11.5, 1.2, BG_CARD)
    add_rect(s, 0.8, y, 0.06, 1.2, color)
    add_text(s, title, 1.1, y+0.1,  7.8, 0.45, size=12, bold=True, color=WHITE)
    add_text(s, desc,  1.1, y+0.62, 7.8, 0.48, size=9,  color=GRAY)
    badge(s, effort, 10.0, y+0.38, 2.1, 0.38, color)


# =============================================================================
# SLIDE 14 - Perguntas
# =============================================================================
s = slide_add("Perguntas")
add_rect(s, 0, 0, 13.33, 7.5, BG_DARK)
add_rect(s, 0, 0, 0.5, 7.5, ACCENT)
add_rect(s, 0.5, 3.6, 12.83, 0.06, ACCENT)

add_text(s, "Obrigado!", 1.2, 0.8, 10.5, 1.5, size=52, bold=True, color=WHITE)
add_text(s, "Perguntas?", 1.2, 2.5, 10.5, 1.0, size=30, color=ACCENT2)

for i, (label, color) in enumerate([("Loki 3.6", ACCENT2), ("Grafana 13.0", ACCENT5),
                                      ("Prometheus 3.11", ACCENT3), ("Alloy 1.9", ACCENT4)]):
    badge(s, label, 1.2 + i * 2.5, 3.85, 2.2, 0.38, color)

add_text(s, "Documentacao: docs/instalacao_airgapped.md  |  docs/playbook_referencia.md  |  docs/dashboards_referencia.md",
         1.2, 4.55, 11.0, 0.5, size=11, color=GRAY)
add_text(s, "Grafana: http://192.168.137.200:3000  |  admin / Obs@2026!",
         1.2, 5.05, 11.0, 0.45, size=11, color=GRAY)


# Salvar
import os
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "stack_observabilidade.pptx")
prs.save(out)
print(f"OK  stack_observabilidade.pptx  ({len(prs.slides)} slides)")
